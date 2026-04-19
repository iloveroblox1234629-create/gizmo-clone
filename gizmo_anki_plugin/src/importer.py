"""Import utilities for various flashcard formats."""
import csv
import re
import subprocess
from pathlib import Path
from typing import List, Dict, Optional

from aqt import mw


def import_quizlet(filepath: str) -> List[Dict[str, str]]:
    """
    Import Quizlet tab-separated export.
    Quizlet export format: Term <tab> Definition
    """
    cards = []
    try:
        with open(filepath, 'r', encoding='utf-8-sig') as f:
            reader = csv.reader(f, delimiter='\t')
            for row in reader:
                if len(row) >= 2:
                    cards.append({
                        "Front": row[0].strip(),
                        "Back": row[1].strip(),
                        "Type": "Quizlet import",
                        "AIHint": ""
                    })
        return cards
    except Exception as e:
        raise ValueError(f"Failed to read Quizlet file: {e}")


def import_anki_csv(filepath: str) -> List[Dict[str, str]]:
    """
    Import Anki-compatible CSV.
    Expected columns: front, back (case-insensitive)
    """
    cards = []
    try:
        with open(filepath, 'r', encoding='utf-8-sig') as f:
            # Detect delimiter
            sample = f.read(1024)
            f.seek(0)
            delimiter = '\t' if '\t' in sample else ','

            reader = csv.DictReader(f, delimiter=delimiter)
            field_map = {}
            for key in reader.fieldnames:
                low = key.lower().strip()
                if low == 'front':
                    field_map['front'] = key
                elif low == 'back':
                    field_map['back'] = key

            if 'front' not in field_map or 'back' not in field_map:
                raise ValueError("CSV must contain 'front' and 'back' columns")

            for row in reader:
                cards.append({
                    "Front": row[field_map['front']].strip(),
                    "Back": row[field_map['back']].strip(),
                    "Type": "CSV import",
                    "AIHint": ""
                })
        return cards
    except Exception as e:
        raise ValueError(f"Failed to read CSV: {e}")


def import_text(filepath: str) -> List[Dict[str, str]]:
    """
    Import plain text file.
    Format: Q: ... A: ... or separated by blank lines.
    """
    cards = []
    try:
        with open(filepath, 'r', encoding='utf-8') as f:
            content = f.read()

        # Split into paragraphs
        paragraphs = [p.strip() for p in content.split('\n\n') if p.strip()]

        if len(paragraphs) >= 2:
            # Consecutive pairs
            for i in range(0, len(paragraphs) - 1, 2):
                cards.append({
                    "Front": paragraphs[i],
                    "Back": paragraphs[i + 1] if i + 1 < len(paragraphs) else "",
                    "Type": "Text import",
                    "AIHint": ""
                })
        return cards
    except Exception as e:
        raise ValueError(f"Failed to read text file: {e}")


def import_pdf(filepath: str) -> List[Dict[str, str]]:
    """
    Import PDF file using subprocess.
    Requires: pdftotext (poppler) or pymupdf
    """
    cards = []

    # Try pymupdf first
    try:
        import fitz  # PyMuPDF
        doc = fitz.open(filepath)
        text = "\n\n".join(page.get_text() for page in doc)
        doc.close()
        return _text_to_cards(text, "PDF (PyMuPDF)")
    except ImportError:
        pass

    # Fall back to pdftotext
    try:
        result = subprocess.run(
            ["pdftotext", filepath, "-"],
            capture_output=True, text=True, timeout=30
        )
        if result.returncode == 0:
            return _text_to_cards(result.stdout, "PDF (pdftotext)")
    except (FileNotFoundError, subprocess.TimeoutExpired):
        pass

    raise ValueError(
        "PDF import requires either PyMuPDF (`pip install pymupdf`) "
        "or pdftotext (poppler-utils)"
    )


def import_pptx(filepath: str) -> List[Dict[str, str]]:
    """
    Import PowerPoint file using python-pptx.
    """
    try:
        from pptx import Presentation
        prs = Presentation(filepath)
        slides_text = []
        for slide in prs.slides:
            slide_text = "\n".join(
                shape.text for shape in slide.shapes
                if hasattr(shape, "text") and shape.text.strip()
            )
            slides_text.append(slide_text)

        text = "\n\n".join(slides_text)
        return _text_to_cards(text, "PowerPoint")
    except ImportError:
        raise ValueError(
            "PowerPoint import requires python-pptx: `pip install python-pptx`"
        )


def import_youtube(url: str) -> List[Dict[str, str]]:
    """
    Fetch YouTube transcript and generate summary cards.
    Requires yt-dlp.
    """
    try:
        import yt_dlp
    except ImportError:
        raise ValueError(
            "YouTube import requires yt-dlp: `pip install yt-dlp`"
        )

    ydl_opts = {
        'writesubtitles': True,
        'subtitleslangs': ['en'],
        'skip_download': True,
        'quiet': True,
    }

    try:
        with yt_dlp.YoutubeDL(ydl_opts) as ydl:
            info = ydl.extract_info(url, download=False)
            title = info.get('title', 'Untitled')

            # Try to get subtitles
            subtitles = info.get('requested_subtitles', {})
            if subtitles:
                sub_url = subtitles.get('en', {}).get('url')
                if sub_url:
                    import requests
                    sub_text = requests.get(sub_url).text
                    # Simple VTT parsing
                    lines = [l.strip() for l in sub_text.split('\n')
                            if l.strip() and not l.startswith(('WEBVTT', 'Kind:', 'Language:'))]
                    # Clean timestamps
                    text = ' '.join(re.sub(r'\d{2}:\d{2}:\d{2}\.\d{3} --> \d{2}:\d{2}:\d{2}\.\d{3}', '', l) for l in lines)
                    text = re.sub(r'<\d+:\d+:\d+\.\d+>', '', text)
                    text = ' '.join(text.split()[:2000])  # Limit

                    return _text_to_cards(text, f"YouTube: {title}")
    except Exception as e:
        raise ValueError(f"YouTube fetch failed: {e}")

    return []


def _text_to_cards(text: str, source: str, max_cards: int = 20) -> List[Dict[str, str]]:
    """
    Parse extracted text into flashcard candidates.
    Heuristic: sentences with keywords, definitions, bold terms, etc.
    """
    cards = []

    # Split into sentences
    sentences = [s.strip() for s in re.split(r'(?<=[.!?])\s+', text) if s.strip()]

    # Find potential question/answer pairs
    # Look for sentences with keywords like "is", "are", "means", "defined as"
    for i, sentence in enumerate(sentences):
        lower = sentence.lower()
        if any(kw in lower for kw in ['is a', 'is an', 'is the', 'are', 'means', 'defined as', 'refers to']):
            # Split into Q & A
            parts = re.split(r'(?:is|are|means|defined as|refers to)\s+', sentence, 1, re.IGNORECASE)
            if len(parts) == 2:
                term = parts[0].strip(' :,-')
                definition = parts[1].strip(' .')
                if term and definition and len(term) < 100:
                    cards.append({
                        "Front": f"What is {term}?",
                        "Back": definition,
                        "Type": f"Auto ({source})",
                        "AIHint": term
                    })

    # Also look for bold/italic terms
    bold_terms = re.findall(r'\*\*(.+?)\*\*|\*(.+?)\*|__(.+?)__|_(.+?)_', text)
    for match in bold_terms:
        term = match[0] or match[1] or match[2] or match[3]
        if term and len(term) < 80:
            cards.append({
                "Front": f"Define: {term}",
                "Back": "",
                "Type": f"Bold term ({source})",
                "AIHint": term
            })

    # Remove duplicates
    seen = set()
    unique = []
    for card in cards:
        key = (card["Front"][:100], card["Back"][:100])
        if key not in seen:
            seen.add(key)
            unique.append(card)

    return unique[:max_cards]


def setup_import_hooks():
    """Register import hooks with Anki's import dialog."""
    from aqt import gui_hooks

    # This would integrate with Anki's import system
    # For now, we use the custom dialog
    pass